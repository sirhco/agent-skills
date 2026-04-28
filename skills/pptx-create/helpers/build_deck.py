"""Reference scaffold for pptx-create skill.

Usage as library:
    from helpers.build_deck import (
        Deck, title_slide, bullets_slide, big_stat_slide,
        chart_slide, image_slide, two_col_slide, table_slide,
        section_slide, quote_slide, cta_slide,
    )
    from themes.themes import get_theme
    d = Deck(theme=get_theme("stripe"))
    d.title("My deck", "Subtitle")
    d.bullets("Agenda", ["A", "B"])
    d.save("out.pptx")

Usage as script:
    python3 build_deck.py corporate   # build sample corporate deck
    python3 build_deck.py pitch       # build sample pitch deck
"""

import os
import sys
from pathlib import Path

# allow running as script or importing from neighbor dir
_HERE = Path(__file__).resolve().parent
_ROOT = _HERE.parent
if str(_ROOT) not in sys.path:
    sys.path.insert(0, str(_ROOT))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from themes.themes import get_theme, CORPORATE_DEFAULT, PITCH_NOIR

__version__ = "0.2.0"


# ---------- primitives ----------

def add_text(slide, text, x, y, w, h, *, size=18, bold=False,
             color=None, align=PP_ALIGN.LEFT, font="Calibri",
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    if color is not None:
        r.font.color.rgb = color
    return tb


def add_bg(slide, prs, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_rect(slide, x, y, w, h, color, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
    return sh


def add_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def add_footer(slide, prs, theme, page_num=None, total=None, brand=""):
    """Bottom strip: brand left, page count right."""
    if brand:
        add_text(slide, brand, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
                 size=10, color=theme["muted"], font=theme["font"])
    if page_num is not None:
        txt = f"{page_num}" if total is None else f"{page_num} / {total}"
        add_text(slide, txt, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
                 size=10, color=theme["muted"], font=theme["font"],
                 align=PP_ALIGN.RIGHT)


# ---------- Deck class ----------

class Deck:
    def __init__(self, theme=None, brand="", brand_kit=None, auto_brand_kit=True):
        """
        theme       : theme dict (from get_theme).
        brand       : footer-left text. Overridden by brand_kit['footer_text'] if set.
        brand_kit   : dict (logo_path, accent, footer_text, font, logo_position).
                      None + auto_brand_kit=True -> load ./brand.toml if present.
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        if brand_kit is None and auto_brand_kit:
            try:
                from helpers.brand_kit import load_brand_kit
                brand_kit = load_brand_kit()
            except Exception:
                brand_kit = {}
        brand_kit = brand_kit or {}
        self.theme = dict(theme) if theme else dict(CORPORATE_DEFAULT)
        # apply brand_kit overrides to theme
        if brand_kit.get("accent"):
            try:
                from themes.themes import _rgb
                self.theme["accent"] = _rgb(brand_kit["accent"])
            except Exception:
                pass
        if brand_kit.get("logo_path"):
            self.theme["logo_path"] = brand_kit["logo_path"]
            self.theme["logo_position"] = brand_kit.get("logo_position", "tr")
        if brand_kit.get("font"):
            self.theme["font"] = brand_kit["font"]
        self.brand = brand_kit.get("footer_text") or brand or self.theme.get("footer_text", "")
        self.brand_kit = brand_kit
        self._slide_index = 0
        self._sections = []   # populated by section() — used for agenda strip
        self._current_section_index = -1
        self._titles = []     # for accessibility uniqueness check

    def _new(self, with_footer=True):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        add_bg(slide, self.prs, self.theme["bg"])
        self._slide_index += 1
        if with_footer and self._slide_index > 1:
            add_footer(slide, self.prs, self.theme,
                       page_num=self._slide_index, brand=self.brand)
            self._draw_logo(slide)
            self._draw_progress(slide)
        return slide

    def _draw_progress(self, slide):
        """Section dots in footer center. Active dot uses accent."""
        if not self._sections:
            return
        t = self.theme
        n = len(self._sections)
        active = self._current_section_index
        # center dots horizontally; small footprint above footer baseline
        slide_w = self.prs.slide_width
        dot_size = Inches(0.10)
        gap = Inches(0.18)
        total_w = n * dot_size + (n - 1) * gap
        x0 = (slide_w - total_w) // 2
        y = Inches(7.10)
        for i in range(n):
            x = x0 + i * (dot_size + gap)
            color = t["accent"] if i == active else t["muted"]
            add_rect(slide, x, y, dot_size, dot_size, color)

    def _draw_logo(self, slide):
        path = self.theme.get("logo_path")
        if not path or not os.path.exists(path):
            return
        pos = self.theme.get("logo_position", "tr")
        # 0.6" square in chosen corner with 0.3" margin
        sw, sh = self.prs.slide_width, self.prs.slide_height
        size = Inches(0.6)
        margin = Inches(0.3)
        if pos == "tl":
            x, y = margin, margin
        elif pos == "tr":
            x, y = sw - size - margin, margin
        elif pos == "bl":
            x, y = margin, sh - size - margin
        else:  # br
            x, y = sw - size - margin, sh - size - margin
        try:
            pic = slide.shapes.add_picture(path, x, y, height=size)
            # accessibility: alt text on logo
            try:
                pic.element.nvSpPr.cNvPr.set("descr", "brand logo")
            except Exception:
                pass
        except Exception:
            pass

    # ---------- slide types ----------

    def title(self, title, subtitle="", notes=""):
        slide = self._new(with_footer=False)
        t = self.theme
        is_pitch = t["mode"] == "pitch"

        if is_pitch:
            add_text(slide, title, Inches(0.6), Inches(2.6), Inches(12), Inches(2.5),
                     size=84, bold=True, color=t["ink"], font=t["font"])
            if subtitle:
                add_text(slide, subtitle, Inches(0.6), Inches(5.5), Inches(12),
                         Inches(0.6), size=20, color=t["muted"], font=t["font"])
            # accent slash
            add_rect(slide, Inches(0.6), Inches(2.4), Inches(0.6), Inches(0.08),
                     t["accent"])
        else:
            add_rect(slide, Inches(0.6), Inches(2.0), Inches(0.6), Inches(0.08),
                     t["accent"])
            add_text(slide, title, Inches(0.6), Inches(2.4), Inches(12), Inches(1.5),
                     size=44, bold=True, color=t["ink"], font=t["font"])
            if subtitle:
                add_text(slide, subtitle, Inches(0.6), Inches(4.0), Inches(12),
                         Inches(0.6), size=20, color=t["muted"], font=t["font"])
        add_notes(slide, notes)
        return slide

    def section(self, label, number=None, notes=""):
        # track section for progress strip on subsequent slides
        self._sections.append(label)
        self._current_section_index = len(self._sections) - 1
        slide = self._new(with_footer=False)
        t = self.theme
        # full-bleed accent panel left third
        add_rect(slide, 0, 0, Inches(4.5), self.prs.slide_height, t["accent"])
        add_text(slide, label, Inches(5.0), Inches(3.0), Inches(8), Inches(2),
                 size=48, bold=True, color=t["ink"], font=t["font"])
        if number is not None:
            num = f"0{number}" if number < 10 else str(number)
            add_text(slide, num, Inches(0.6), Inches(0.6), Inches(2), Inches(0.6),
                     size=14, color=t["bg"], font=t["font"])
        add_notes(slide, notes)
        return slide

    def bullets(self, title, bullets, *, icons=None, notes=""):
        """Bullets slide. Auto-shrinks font 20pt -> 14pt floor; splits to continuation slide if still over.

        icons       : optional list of icon names (resolved against helpers/icons/) parallel to bullets.
        """
        # estimate space: each bullet at size pt occupies ~ (size * 1.4) pt + 14pt space-after
        AVAIL_HEIGHT_PT = 5 * 72  # 5 inches usable
        FONT_SIZES = [20, 18, 16, 15, 14]
        SPACE_AFTER = 14

        def fits(items, size):
            line_h = size * 1.4 + SPACE_AFTER
            # also penalize long bullets that wrap: rough char/line cap
            chars_per_line = max(40, int(120 - size * 2))
            total = 0
            for b in items:
                wraps = max(1, (len(b) // chars_per_line) + (1 if len(b) % chars_per_line else 0))
                total += wraps * line_h
            return total <= AVAIL_HEIGHT_PT

        # find smallest size that fits all bullets, else split
        chosen_size = None
        for size in FONT_SIZES:
            if fits(bullets, size):
                chosen_size = size
                break

        if chosen_size is None:
            # split: largest prefix that fits at floor size, recurse for rest
            floor = FONT_SIZES[-1]
            prefix = []
            for b in bullets:
                trial = prefix + [b]
                if fits(trial, floor):
                    prefix = trial
                else:
                    break
            if not prefix:
                # single bullet too long even at floor — render it alone, truncated visually by overflow
                prefix = [bullets[0]]
            rest = bullets[len(prefix):]
            self._render_bullets_slide(title, prefix, icons=icons, size=floor, notes=notes)
            if rest:
                rest_icons = icons[len(prefix):] if icons else None
                self.bullets(f"{title} (cont.)", rest, icons=rest_icons, notes="")
            return None

        return self._render_bullets_slide(title, bullets, icons=icons, size=chosen_size, notes=notes)

    def _render_bullets_slide(self, title, bullets, *, icons=None, size=20, notes=""):
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)

        x = Inches(0.6)
        y = Inches(1.8)
        w = Inches(12)
        h = Inches(5)

        if icons:
            # resolve icons; left-pad bullets with icon column
            icon_w = Inches(0.4)
            icon_pad = Inches(0.15)
            text_x = x + icon_w + icon_pad
            text_w = w - icon_w - icon_pad
            row_h = Pt(size * 1.6 + 14) // 1
            from helpers.icons import resolve_icon
            for i, b in enumerate(bullets):
                cy = y + i * row_h
                ipath = resolve_icon(icons[i]) if i < len(icons) and icons[i] else None
                if ipath:
                    try:
                        slide.shapes.add_picture(str(ipath), x, cy, width=icon_w, height=icon_w)
                    except Exception:
                        pass
                tb = slide.shapes.add_textbox(text_x, cy, text_w, row_h)
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = b
                r.font.name = t["font"]
                r.font.size = Pt(size)
                r.font.color.rgb = t["ink"]
        else:
            tb = slide.shapes.add_textbox(x, y, w, h)
            tf = tb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(14)
                r = p.add_run()
                r.text = f"•  {b}"
                r.font.name = t["font"]
                r.font.size = Pt(size)
                r.font.color.rgb = t["ink"]
        add_notes(slide, notes)
        return slide

    def two_col(self, title, left_title, left_items, right_title, right_items,
                notes=""):
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)

        for col, (h, items) in enumerate([(left_title, left_items),
                                          (right_title, right_items)]):
            x = Inches(0.6 + col * 6.2)
            add_text(slide, h, x, Inches(1.9), Inches(5.8), Inches(0.6),
                     size=16, bold=True, color=t["accent"], font=t["font"])
            tb = slide.shapes.add_textbox(x, Inches(2.6), Inches(5.8), Inches(4.5))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(10)
                r = p.add_run()
                r.text = f"•  {b}"
                r.font.name = t["font"]
                r.font.size = Pt(18)
                r.font.color.rgb = t["ink"]
        add_notes(slide, notes)
        return slide

    def big_stat(self, stat, caption, notes=""):
        slide = self._new(with_footer=False)
        t = self.theme
        add_text(slide, stat, Inches(0.6), Inches(1.5), Inches(12), Inches(3.5),
                 size=200, bold=True, color=t["accent"], font=t["font"])
        add_text(slide, caption, Inches(0.6), Inches(5.6), Inches(12), Inches(1),
                 size=22, color=t["ink"], font=t["font"])
        add_notes(slide, notes)
        return slide

    def quote(self, quote, attribution="", notes=""):
        slide = self._new()
        t = self.theme
        add_text(slide, f"“{quote}”", Inches(1.5), Inches(2.0),
                 Inches(10.3), Inches(3.5), size=32, bold=False,
                 color=t["ink"], font=t["font"])
        if attribution:
            add_text(slide, f"— {attribution}", Inches(1.5), Inches(5.5),
                     Inches(10.3), Inches(0.6),
                     size=16, color=t["muted"], font=t["font"])
        add_notes(slide, notes)
        return slide

    _CHART_KINDS = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
        "line":   XL_CHART_TYPE.LINE,
        "pie":    XL_CHART_TYPE.PIE,
        "area":   XL_CHART_TYPE.AREA,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }

    def chart(self, title, categories, series, *, kind="column",
              chart_type=None, notes=""):
        """Native editable PPTX chart.

        kind        : 'column' (default), 'bar', 'line', 'pie', 'area', 'doughnut'.
        chart_type  : raw XL_CHART_TYPE override (advanced).
        series      : dict[name -> iterable of values].
        """
        slide = self._new()
        self._title_bar(slide, title)
        ct = chart_type or self._CHART_KINDS.get(kind, XL_CHART_TYPE.COLUMN_CLUSTERED)
        data = CategoryChartData()
        data.categories = categories
        for name, values in series.items():
            data.add_series(name, values)
        gframe = slide.shapes.add_chart(
            ct, Inches(0.6), Inches(1.8), Inches(12), Inches(5), data
        )
        # subtle styling: hide chart title (we render slide title separately)
        try:
            ch = gframe.chart
            ch.has_title = False
            if len(series) > 1:
                ch.has_legend = True
                ch.legend.include_in_layout = False
            else:
                ch.has_legend = False
        except Exception:
            pass
        add_notes(slide, notes)
        return slide

    def table(self, title, headers, rows, notes=""):
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)
        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl = slide.shapes.add_table(
            n_rows, n_cols, Inches(0.6), Inches(1.9), Inches(12), Inches(5)
        ).table
        for c, h in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = h
            r.font.name = t["font"]
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = t["bg"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = t["accent"]
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = ""
                tf = cell.text_frame
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = str(val)
                r.font.name = t["font"]
                r.font.size = Pt(13)
                r.font.color.rgb = t["ink"]
        add_notes(slide, notes)
        return slide

    def image(self, title, image_path, alt=None, caption="", *, fit="contain", notes=""):
        """Image slide.

        alt         : REQUIRED. Screen-reader alt text. Pass empty string only for decorative images.
        fit         : 'contain' (letterbox, default) or 'cover' (crop to fill).
        caption     : small line below the image.
        """
        if alt is None:
            raise TypeError("Deck.image() missing required keyword argument: 'alt' (use alt='' for decorative)")
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)

        # available area
        ax, ay, aw, ah = Inches(0.6), Inches(1.9), Inches(12), Inches(4.7)
        if image_path and os.path.exists(image_path):
            try:
                pic = self._add_fitted_picture(slide, image_path, ax, ay, aw, ah, fit=fit)
                # accessibility — set descr on picture
                try:
                    pic.element.nvSpPr.cNvPr.set("descr", alt or "")
                    if alt:
                        pic.element.nvSpPr.cNvPr.set("title", alt[:60])
                except Exception:
                    pass
            except Exception as e:
                add_text(slide, f"[image error: {e}]",
                         Inches(0.6), Inches(3.5), Inches(12), Inches(0.6),
                         size=14, color=t["muted"], font=t["font"])
        else:
            add_text(slide, f"[image not found: {image_path}]",
                     Inches(0.6), Inches(3.5), Inches(12), Inches(0.6),
                     size=14, color=t["muted"], font=t["font"])
        if caption:
            add_text(slide, caption, Inches(0.6), Inches(6.7), Inches(12),
                     Inches(0.4), size=12, color=t["muted"], font=t["font"])
        add_notes(slide, notes)
        return slide

    def _add_fitted_picture(self, slide, image_path, ax, ay, aw, ah, *, fit="contain"):
        """Add picture honoring fit=contain|cover within (ax,ay,aw,ah)."""
        # measure image aspect via Pillow if available; else letterbox by width
        try:
            from PIL import Image
            with Image.open(image_path) as im:
                iw, ih = im.size
        except Exception:
            return slide.shapes.add_picture(image_path, ax, ay, width=aw)

        ar_img = iw / ih
        ar_box = aw / ah
        if fit == "cover":
            # fill box; crop excess implicitly by sizing past edges then re-clipping via box
            if ar_img > ar_box:
                # image wider than box: match height
                new_h = ah
                new_w = int(ah * ar_img)
                x = ax - (new_w - aw) // 2
                y = ay
            else:
                new_w = aw
                new_h = int(aw / ar_img)
                x = ax
                y = ay - (new_h - ah) // 2
            return slide.shapes.add_picture(image_path, x, y, width=new_w, height=new_h)
        # contain (default): fit inside box, center
        if ar_img > ar_box:
            new_w = aw
            new_h = int(aw / ar_img)
            x = ax
            y = ay + (ah - new_h) // 2
        else:
            new_h = ah
            new_w = int(ah * ar_img)
            x = ax + (aw - new_w) // 2
            y = ay
        return slide.shapes.add_picture(image_path, x, y, width=new_w, height=new_h)

    def cta(self, headline, sub="", contact="", notes=""):
        slide = self._new(with_footer=False)
        t = self.theme
        add_rect(slide, 0, 0, self.prs.slide_width, self.prs.slide_height,
                 t["accent"])
        add_text(slide, headline, Inches(0.6), Inches(2.6), Inches(12),
                 Inches(2), size=64, bold=True, color=t["bg"],
                 font=t["font"])
        if sub:
            add_text(slide, sub, Inches(0.6), Inches(4.8), Inches(12),
                     Inches(0.8), size=24, color=t["bg"], font=t["font"])
        if contact:
            add_text(slide, contact, Inches(0.6), Inches(6.4), Inches(12),
                     Inches(0.6), size=16, color=t["bg"], font=t["font"])
        add_notes(slide, notes)
        return slide

    # ---------- additional slide kinds ----------

    def timeline(self, title, events, notes=""):
        """Horizontal timeline. events: list of (date_str, label_str)."""
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)
        if not events:
            return slide
        n = len(events)
        x_left = Inches(0.8)
        x_right = self.prs.slide_width - Inches(0.8)
        y_axis = Inches(4.0)
        # axis line
        add_rect(slide, x_left, y_axis, x_right - x_left, Inches(0.04), t["accent"])
        step = (x_right - x_left) // max(n - 1, 1)
        for i, (date, label) in enumerate(events):
            cx = x_left + i * step if n > 1 else (x_left + x_right) // 2
            # node dot
            dot = Inches(0.22)
            add_rect(slide, cx - dot // 2, y_axis - dot // 2 + Inches(0.02), dot, dot, t["accent"])
            # alternate label above/below
            above = (i % 2 == 0)
            label_y = y_axis - Inches(1.4) if above else y_axis + Inches(0.5)
            date_y = label_y + (Inches(0.0) if not above else Inches(0.6))
            add_text(slide, label, cx - Inches(1.2), label_y, Inches(2.4), Inches(0.6),
                     size=14, color=t["ink"], font=t["font"], align=PP_ALIGN.CENTER)
            add_text(slide, date, cx - Inches(1.2), date_y, Inches(2.4), Inches(0.4),
                     size=11, color=t["muted"], font=t["font"], align=PP_ALIGN.CENTER)
        add_notes(slide, notes)
        return slide

    def comparison(self, title, headers, rows, notes=""):
        """Comparison matrix. headers: ['Feature', 'Us', 'Them'] (first col is row label).
        rows: list of [label, val_a, val_b, ...]. Truthy strings render as ✓."""
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)
        # render as table; replace truthy "yes"/"true"/"✓" with bold check
        n_rows = len(rows) + 1
        n_cols = len(headers)
        tbl = slide.shapes.add_table(
            n_rows, n_cols, Inches(0.6), Inches(1.9), Inches(12), Inches(5)
        ).table
        for c, h in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = ""
            tf = cell.text_frame
            r = tf.paragraphs[0].add_run()
            r.text = h
            r.font.name = t["font"]
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = t["bg"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = t["accent"]
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = ""
                tf = cell.text_frame
                p = tf.paragraphs[0]
                r = p.add_run()
                s = str(val).strip()
                lc = s.lower()
                if ci > 0 and lc in ("yes", "true", "y", "✓", "1"):
                    r.text = "✓"
                    r.font.bold = True
                    r.font.color.rgb = t["accent"]
                elif ci > 0 and lc in ("no", "false", "n", "✗", "0", "—", "-"):
                    r.text = "—"
                    r.font.color.rgb = t["muted"]
                else:
                    r.text = s
                    r.font.color.rgb = t["ink"]
                r.font.name = t["font"]
                r.font.size = Pt(14)
        add_notes(slide, notes)
        return slide

    def metric_grid(self, title, metrics, notes=""):
        """4-up KPI tiles. metrics: list of (value, label, delta) — delta optional, signed string."""
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)
        n = min(len(metrics), 4)
        if n == 0:
            return slide
        # 2x2 (or 1xN if <=2) grid
        cols = 2 if n > 2 else n
        rows = (n + cols - 1) // cols
        gx, gy = Inches(0.6), Inches(1.9)
        gw, gh = Inches(12), Inches(5)
        cell_w = (gw - Inches(0.4) * (cols - 1)) // cols
        cell_h = (gh - Inches(0.4) * (rows - 1)) // rows
        for i in range(n):
            r, c = divmod(i, cols)
            x = gx + c * (cell_w + Inches(0.4))
            y = gy + r * (cell_h + Inches(0.4))
            add_rect(slide, x, y, cell_w, cell_h, t["accent2"])
            value, label = metrics[i][0], metrics[i][1]
            delta = metrics[i][2] if len(metrics[i]) > 2 else ""
            add_text(slide, str(value), x + Inches(0.3), y + Inches(0.3), cell_w - Inches(0.6),
                     Inches(1.4), size=54, bold=True, color=t["accent"], font=t["font"])
            add_text(slide, str(label), x + Inches(0.3), y + Inches(1.8), cell_w - Inches(0.6),
                     Inches(0.5), size=14, color=t["ink"], font=t["font"])
            if delta:
                color = t["accent"] if str(delta).startswith("+") else t["muted"]
                add_text(slide, str(delta), x + Inches(0.3), y + Inches(2.3), cell_w - Inches(0.6),
                         Inches(0.4), size=12, color=color, font=t["font"])
        add_notes(slide, notes)
        return slide

    def code(self, title, code_text, lang="python", notes=""):
        """Code slide. Mono font; pygments syntax highlight if installed (best-effort)."""
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)
        # background panel
        add_rect(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(5), t["accent2"])

        tb = slide.shapes.add_textbox(Inches(0.85), Inches(2.05), Inches(11.5), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = False

        try:
            from pygments import lex
            from pygments.lexers import get_lexer_by_name, guess_lexer
            try:
                lexer = get_lexer_by_name(lang)
            except Exception:
                lexer = guess_lexer(code_text)
            tokens = list(lex(code_text, lexer))
            self._render_pygmented(tf, tokens, t)
        except ImportError:
            # plain mono
            for i, line in enumerate(code_text.splitlines() or [""]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p.add_run()
                r.text = line
                r.font.name = "Consolas"
                r.font.size = Pt(14)
                r.font.color.rgb = t["ink"]
        add_notes(slide, notes)
        return slide

    def _render_pygmented(self, tf, tokens, theme):
        """Render pygments tokens into a text_frame with simple color mapping."""
        from pygments.token import Token
        # crude palette — accent for keywords, muted for comments, ink for default
        def color_for(tok):
            if tok in Token.Comment:
                return theme["muted"]
            if tok in Token.Keyword:
                return theme["accent"]
            if tok in Token.String:
                return theme["accent"]
            if tok in Token.Name.Function or tok in Token.Name.Class:
                return theme["accent"]
            if tok in Token.Number:
                return theme["accent"]
            return theme["ink"]

        first_line = True
        cur_p = tf.paragraphs[0]
        for tok, text in tokens:
            if "\n" in text:
                parts = text.split("\n")
                for i, part in enumerate(parts):
                    if part:
                        run = cur_p.add_run()
                        run.text = part
                        run.font.name = "Consolas"
                        run.font.size = Pt(14)
                        run.font.color.rgb = color_for(tok)
                    if i < len(parts) - 1:
                        cur_p = tf.add_paragraph()
                continue
            if not text:
                continue
            run = cur_p.add_run()
            run.text = text
            run.font.name = "Consolas"
            run.font.size = Pt(14)
            run.font.color.rgb = color_for(tok)

    def agenda(self, title, items, current=None, notes=""):
        """Agenda / TOC. items: list of strings. current: index to highlight."""
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(12), Inches(5))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(12)
            num = p.add_run()
            num.text = f"{i+1:02d}   "
            num.font.name = t["font"]
            num.font.size = Pt(20)
            num.font.color.rgb = t["accent"] if i == current else t["muted"]
            num.font.bold = (i == current)
            r = p.add_run()
            r.text = item
            r.font.name = t["font"]
            r.font.size = Pt(20)
            r.font.color.rgb = t["ink"] if i == current else t["muted"]
            r.font.bold = (i == current)
        add_notes(slide, notes)
        return slide

    def pricing(self, title, plans, notes=""):
        """Pricing tiers. plans: list of dicts {name, price, period, features:[...], featured:bool}."""
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)
        n = max(1, min(len(plans), 4))
        gx, gy = Inches(0.6), Inches(1.9)
        gw, gh = Inches(12), Inches(5)
        col_w = (gw - Inches(0.3) * (n - 1)) // n
        for i, plan in enumerate(plans[:n]):
            x = gx + i * (col_w + Inches(0.3))
            y = gy
            featured = bool(plan.get("featured"))
            bg = t["accent"] if featured else t["accent2"]
            ink = t["bg"] if featured else t["ink"]
            sub = t["bg"] if featured else t["muted"]
            add_rect(slide, x, y, col_w, gh, bg)
            add_text(slide, plan.get("name", ""), x + Inches(0.3), y + Inches(0.3),
                     col_w - Inches(0.6), Inches(0.5), size=18, bold=True, color=ink, font=t["font"])
            price = str(plan.get("price", ""))
            period = plan.get("period", "")
            add_text(slide, price, x + Inches(0.3), y + Inches(0.9), col_w - Inches(0.6),
                     Inches(1.0), size=40, bold=True, color=ink, font=t["font"])
            if period:
                add_text(slide, period, x + Inches(0.3), y + Inches(1.95), col_w - Inches(0.6),
                         Inches(0.4), size=12, color=sub, font=t["font"])
            features = plan.get("features", [])[:6]
            tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(2.5),
                                          col_w - Inches(0.6), gh - Inches(2.7))
            tf = tb.text_frame
            tf.word_wrap = True
            for j, feat in enumerate(features):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_after = Pt(6)
                r = p.add_run()
                r.text = f"✓  {feat}"
                r.font.name = t["font"]
                r.font.size = Pt(12)
                r.font.color.rgb = ink
        add_notes(slide, notes)
        return slide

    def before_after(self, title, before, after, notes=""):
        """Two-column before/after. before/after: dict {label, body|image_path|alt}."""
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)
        for col, (kind, data) in enumerate([("BEFORE", before), ("AFTER", after)]):
            x = Inches(0.6 + col * 6.2)
            label_color = t["muted"] if col == 0 else t["accent"]
            add_text(slide, kind, x, Inches(1.9), Inches(5.8), Inches(0.4),
                     size=12, bold=True, color=label_color, font=t["font"])
            if data.get("image_path") and os.path.exists(data["image_path"]):
                self._add_fitted_picture(slide, data["image_path"],
                                         x, Inches(2.4), Inches(5.8), Inches(3.5),
                                         fit=data.get("fit", "contain"))
            else:
                add_text(slide, data.get("body", ""), x, Inches(2.4), Inches(5.8), Inches(3.5),
                         size=18, color=t["ink"], font=t["font"])
            if data.get("label"):
                add_text(slide, data["label"], x, Inches(6.0), Inches(5.8), Inches(0.5),
                         size=14, color=t["muted"], font=t["font"])
        add_notes(slide, notes)
        return slide

    def video_placeholder(self, title, alt, poster_path=None, notes=""):
        """Frame + play icon. No actual video embed (avoids huge file size). alt required."""
        if not alt:
            raise TypeError("Deck.video_placeholder() requires alt= for accessibility.")
        slide = self._new()
        t = self.theme
        self._title_bar(slide, title)
        ax, ay, aw, ah = Inches(0.6), Inches(1.9), Inches(12), Inches(4.7)
        if poster_path and os.path.exists(poster_path):
            self._add_fitted_picture(slide, poster_path, ax, ay, aw, ah, fit="cover")
        else:
            add_rect(slide, ax, ay, aw, ah, t["accent2"])
        # play triangle
        cx = ax + aw // 2
        cy = ay + ah // 2
        size = Inches(0.8)
        add_text(slide, "▶", cx - size, cy - size, size * 2, size * 2,
                 size=120, color=t["accent"], font=t["font"], align=PP_ALIGN.CENTER)
        add_text(slide, alt, ax, ay + ah + Inches(0.1), aw, Inches(0.4),
                 size=12, color=t["muted"], font=t["font"], align=PP_ALIGN.CENTER)
        add_notes(slide, notes)
        return slide

    # ---------- internal ----------

    def _title_bar(self, slide, title):
        t = self.theme
        add_text(slide, title, Inches(0.6), Inches(0.6), Inches(12), Inches(0.9),
                 size=30, bold=True, color=t["ink"], font=t["font"])
        add_rect(slide, Inches(0.6), Inches(1.45), Inches(0.6), Inches(0.06),
                 t["accent"])

    # ---------- output ----------

    def save(self, path):
        self.prs.save(path)
        return path


# ---------- sample builds (run as script) ----------

def _sample_corporate(out="deck_corporate.pptx"):
    d = Deck(theme=get_theme("stripe"), brand="Acme Corp · 2026")
    d.title("Q3 Business Review",
            "Operations + Product · 2026-04-27",
            notes="Open with topline. 90s.")
    d.bullets("Agenda",
              ["Topline metrics", "Product highlights", "Risks", "Q4 plan"])
    d.section("Topline", number=1)
    d.chart("Revenue by quarter",
            ["Q1", "Q2", "Q3", "Q4 (proj)"],
            {"Revenue ($M)": (12.4, 15.1, 18.7, 22.0)},
            notes="Q3 beat plan by 8%.")
    d.table("Customer mix",
            ["Segment", "ARR ($M)", "Growth"],
            [["Enterprise", "11.2", "+34%"],
             ["Mid-market", "5.8", "+22%"],
             ["SMB", "1.7", "+11%"]])
    d.two_col("Risks vs Mitigations",
              "Risks", ["Hiring lag in EMEA", "One large customer at renewal",
                        "Infra cost trending up 12%"],
              "Mitigations", ["Open EMEA contractor pool", "Exec sponsor on account",
                              "Reserved capacity contracts"])
    d.quote("We don't ship to schedule. We ship when it's right.",
            attribution="CEO, all-hands 2026-Q1")
    d.cta("Q4 plan: ship billing v2",
          sub="Close 3 enterprise deals · Hire VP Eng · Reduce infra 8%")
    return d.save(out)


def _sample_pitch(out="deck_pitch.pptx"):
    d = Deck(theme=get_theme("pitch-noir"), brand="Ferro")
    d.title("Ferro", "AI dev tools that ship themselves.")
    d.big_stat("73%",
               "of engineering time is spent reading code, not writing it.")
    d.bullets("What we built",
              ["Autonomous PR review",
               "Zero-config setup",
               "Works with any stack"])
    d.big_stat("$1.2M", "ARR · 4 months · 38 customers")
    d.chart("Growth", ["M1", "M2", "M3", "M4"],
            {"ARR ($K)": (180, 420, 780, 1200)})
    d.cta("Raising $8M", sub="Seed · close 2026-06", contact="chris@ferro.ai")
    return d.save(out)


if __name__ == "__main__":
    mode = sys.argv[1] if len(sys.argv) > 1 else "corporate"
    path = _sample_pitch() if mode == "pitch" else _sample_corporate()
    prs = Presentation(path)
    print(f"saved: {path}")
    print(f"slides: {len(prs.slides)}")
    print(f"size:   {os.path.getsize(path)} bytes")

"""PDF + thumbnails + reveal.js export."""
from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path


def _libreoffice_bin() -> str | None:
    for name in ("soffice", "libreoffice"):
        bin_ = shutil.which(name)
        if bin_:
            return bin_
    # macOS app bundle common path
    macos = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if Path(macos).is_file():
        return macos
    return None


def export_pdf(deck: Path, out: Path) -> int:
    bin_ = _libreoffice_bin()
    if not bin_:
        print("error: LibreOffice not found. install from https://libreoffice.org or "
              "use --thumbs for PNGs.", file=sys.stderr)
        return 1
    out.parent.mkdir(parents=True, exist_ok=True)
    tmp_dir = out.parent / "_pdf_tmp"
    tmp_dir.mkdir(exist_ok=True)
    cmd = [bin_, "--headless", "--convert-to", "pdf",
           "--outdir", str(tmp_dir), str(deck)]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"error: libreoffice failed: {result.stderr}", file=sys.stderr)
        return 1
    produced = tmp_dir / (deck.stem + ".pdf")
    if not produced.is_file():
        print(f"error: pdf not produced at {produced}", file=sys.stderr)
        return 1
    produced.replace(out)
    try:
        tmp_dir.rmdir()
    except OSError:
        pass
    print(f"+ {out}")
    return 0


def export_thumbs(deck: Path, out_dir: Path) -> int:
    """Render each slide to PNG via libreoffice + pdf -> images. Falls back to placeholder."""
    bin_ = _libreoffice_bin()
    if not bin_:
        print("error: LibreOffice not found; thumbnails require it.", file=sys.stderr)
        return 1
    out_dir.mkdir(parents=True, exist_ok=True)
    pdf = out_dir / (deck.stem + ".pdf")
    rc = export_pdf(deck, pdf)
    if rc != 0:
        return rc
    # PDF -> PNGs
    try:
        from pdf2image import convert_from_path  # type: ignore
        pages = convert_from_path(str(pdf), dpi=120)
        for i, p in enumerate(pages, start=1):
            p.save(out_dir / f"slide_{i:03d}.png", "PNG")
        pdf.unlink(missing_ok=True)
        print(f"+ {len(pages)} thumbs in {out_dir}")
        return 0
    except ImportError:
        # fallback: try pdftoppm (poppler) if installed
        if shutil.which("pdftoppm"):
            base = out_dir / "slide"
            subprocess.run(["pdftoppm", "-png", "-r", "120", str(pdf), str(base)], check=True)
            pdf.unlink(missing_ok=True)
            print(f"+ thumbs in {out_dir} (pdftoppm)")
            return 0
        print("error: install pdf2image (pip) or poppler (pdftoppm) to render thumbnails.",
              file=sys.stderr)
        return 1


def export_reveal(deck: Path, out_dir: Path) -> int:
    """Render the deck's text content as a minimal Reveal.js HTML bundle.

    This is text-only — images/charts/tables are not transferred since the source
    is a binary .pptx. Use markdown_to_pptx + a separate md->reveal pipeline for
    full fidelity. Provided here as a quick-share fallback.
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    from pptx import Presentation
    prs = Presentation(str(deck))
    sections = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        body = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if not title:
                    title = text.splitlines()[0]
                    body.extend(text.splitlines()[1:])
                else:
                    body.extend(text.splitlines())
        body_html = ""
        if body:
            body_html = "<ul>" + "".join(f"<li>{_esc(line)}</li>" for line in body if line.strip()) + "</ul>"
        sections.append(f"<section><h2>{_esc(title) or f'Slide {idx}'}</h2>{body_html}</section>")
    html = _REVEAL_TEMPLATE.replace("__TITLE__", _esc(deck.stem)).replace("__SECTIONS__", "\n".join(sections))
    (out_dir / "index.html").write_text(html, encoding="utf-8")
    print(f"+ {out_dir / 'index.html'} (reveal.js text-only)")
    return 0


def _esc(s: str) -> str:
    return (s or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


_REVEAL_TEMPLATE = """<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<title>__TITLE__</title>
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5/dist/reset.css">
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5/dist/reveal.css">
<link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@5/dist/theme/white.css">
</head>
<body>
<div class="reveal"><div class="slides">
__SECTIONS__
</div></div>
<script src="https://cdn.jsdelivr.net/npm/reveal.js@5/dist/reveal.js"></script>
<script>Reveal.initialize({hash:true});</script>
</body>
</html>
"""

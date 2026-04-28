"""Inspect a .pptx deck — pure-Python, no external binaries.

Reports per-slide structure: title, text, shape counts, charts, tables, images,
speaker notes. Use this to verify a generated deck without opening PowerPoint.

Run:
  python3 inspect_deck.py deck.pptx                  # markdown report to stdout
  python3 inspect_deck.py deck.pptx --json           # JSON report
  python3 inspect_deck.py deck.pptx --notes          # include speaker notes
  python3 inspect_deck.py deck.pptx --render-pdf     # optional, requires LibreOffice
                                                     # (skip on corporate machines)
"""

import argparse
import json
import os
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _shape_summary(shape):
    if shape.has_text_frame:
        txt = shape.text_frame.text.strip()
        return ("text", txt[:120]) if txt else None
    if shape.has_table:
        t = shape.table
        return ("table", f"{len(t.rows)}×{len(t.columns)}")
    if shape.has_chart:
        return ("chart", shape.chart.chart_type)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return ("image", "")
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return ("shape", "")
    return None


def inspect(path, include_notes=False):
    p = Presentation(path)
    width_in = p.slide_width / 914400
    height_in = p.slide_height / 914400

    out = {
        "path": str(Path(path).resolve()),
        "size_bytes": os.path.getsize(path),
        "slide_count": len(p.slides),
        "dimensions_in": [round(width_in, 2), round(height_in, 2)],
        "slides": [],
    }

    for i, slide in enumerate(p.slides, 1):
        items = []
        title = None
        for sh in slide.shapes:
            s = _shape_summary(sh)
            if s is None:
                continue
            items.append({"kind": s[0], "info": str(s[1])})
            if s[0] == "text" and title is None and len(s[1]) <= 120:
                title = s[1].splitlines()[0] if s[1] else None
        notes = ""
        if include_notes and slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        out["slides"].append({
            "n": i,
            "title": title or "(untitled)",
            "shapes": items,
            "notes": notes,
        })
    return out


def to_markdown(report, include_notes=False):
    lines = []
    lines.append(f"# Deck inspection — {Path(report['path']).name}")
    lines.append("")
    lines.append(f"- path: `{report['path']}`")
    lines.append(f"- size: {report['size_bytes']:,} bytes")
    lines.append(f"- slides: {report['slide_count']}")
    lines.append(f"- dimensions: {report['dimensions_in'][0]} × "
                 f"{report['dimensions_in'][1]} in")
    lines.append("")
    for s in report["slides"]:
        lines.append(f"## {s['n']}. {s['title']}")
        kinds = {}
        for sh in s["shapes"]:
            kinds[sh["kind"]] = kinds.get(sh["kind"], 0) + 1
        if kinds:
            lines.append("- shapes: " + ", ".join(
                f"{k}×{v}" for k, v in sorted(kinds.items())))
        # show first text shape body for readability
        for sh in s["shapes"]:
            if sh["kind"] == "text" and sh["info"] != s["title"]:
                lines.append(f"  - text: {sh['info']}")
        if include_notes and s["notes"]:
            lines.append(f"- notes: {s['notes']}")
        lines.append("")
    return "\n".join(lines)


def render_pdf_optional(pptx_path, out_dir=None):
    """Optional: convert to PDF via LibreOffice, ONLY if installed.

    On corporate machines without LibreOffice, this is silently skipped.
    Returns path to PDF or None.
    """
    soffice = None
    for cand in ("libreoffice", "soffice",
                 "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        if shutil.which(cand) or Path(cand).exists():
            soffice = cand
            break
    if not soffice:
        print("LibreOffice not found — skipping PDF render. "
              "(This is fine — inspection report uses python-pptx only.)",
              file=sys.stderr)
        return None
    pptx_path = Path(pptx_path).resolve()
    out_dir = Path(out_dir or pptx_path.parent)
    cmd = [soffice, "--headless", "--convert-to", "pdf",
           "--outdir", str(out_dir), str(pptx_path)]
    subprocess.run(cmd, check=True)
    return out_dir / (pptx_path.stem + ".pdf")


def main(argv=None):
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--json", action="store_true", help="JSON output")
    ap.add_argument("--notes", action="store_true", help="include speaker notes")
    ap.add_argument("--render-pdf", action="store_true",
                    help="optional PDF via LibreOffice — skipped if not installed")
    args = ap.parse_args(argv)

    report = inspect(args.pptx, include_notes=args.notes)
    if args.json:
        print(json.dumps(report, indent=2))
    else:
        print(to_markdown(report, include_notes=args.notes))

    if args.render_pdf:
        pdf = render_pdf_optional(args.pptx)
        if pdf:
            print(f"\nPDF: {pdf}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())

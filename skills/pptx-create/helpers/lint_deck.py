"""Lint a .pptx for accessibility / quality issues.

Checks:
- E: image without alt text (descr attribute empty/missing)
- E: WCAG AA contrast violation (text vs background, ratio < 4.5)
- W: text shape likely overflows (line count > capacity for box)
- W: duplicate slide titles
- W: inconsistent title capitalization (mix of title-case + sentence-case)

Returns list[Issue]. CLI prints + exits nonzero on errors.
"""
from __future__ import annotations

import re
from dataclasses import dataclass, asdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


@dataclass
class Issue:
    severity: str   # "error" | "warning"
    code: str
    slide: int
    message: str

    def to_dict(self):
        return asdict(self)


def _luminance(rgb_hex: str) -> float:
    """sRGB relative luminance per WCAG."""
    h = rgb_hex.lstrip("#")
    if len(h) != 6:
        return 0.0
    r, g, b = (int(h[i:i+2], 16) / 255.0 for i in (0, 2, 4))
    def chan(c):
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    R, G, B = chan(r), chan(g), chan(b)
    return 0.2126 * R + 0.7152 * G + 0.0722 * B


def _contrast(a: str, b: str) -> float:
    la = _luminance(a) + 0.05
    lb = _luminance(b) + 0.05
    return la / lb if la >= lb else lb / la


def _color_hex(color_obj) -> str | None:
    """Extract hex from python-pptx ColorFormat. None if unresolvable (theme color)."""
    try:
        rgb = color_obj.rgb
        if rgb is None:
            return None
        return "#" + str(rgb).upper()
    except Exception:
        return None


def _slide_bg_hex(slide) -> str | None:
    """Best-effort: look for full-bleed rectangle as bg, else None."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            continue
        if not getattr(shape, "fill", None):
            continue
        try:
            if shape.fill.type == 1:  # MSO_FILL.SOLID
                return _color_hex(shape.fill.fore_color)
        except Exception:
            continue
    return None


def lint(deck_path: Path) -> list[Issue]:
    issues: list[Issue] = []
    prs = Presentation(str(deck_path))

    titles = []
    for idx, slide in enumerate(prs.slides, start=1):
        # gather title (first non-empty paragraph that looks like a title)
        slide_title = _detect_title(slide)
        if slide_title:
            titles.append((idx, slide_title))

        # image alt-text check
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                descr = shape.element.nvSpPr.cNvPr.get("descr") or ""
                title_attr = shape.element.nvSpPr.cNvPr.get("title") or ""
                # decorative images may legitimately have empty alt — flag empty as warning, missing as error
                if not descr.strip() and not title_attr.strip():
                    # tolerate logos (already set descr=brand logo upstream); look for that hint
                    issues.append(Issue("error", "alt-missing", idx,
                                        f"picture has no alt text (descr/title unset)"))

        # contrast check on text vs slide bg
        bg_hex = _slide_bg_hex(slide)
        if bg_hex:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        fc = _color_hex(run.font.color)
                        if fc is None:
                            continue
                        ratio = _contrast(fc, bg_hex)
                        if ratio < 4.5:
                            issues.append(Issue(
                                "error", "contrast",
                                idx,
                                f"text {fc} on bg {bg_hex} contrast {ratio:.2f} < 4.5 (WCAG AA): {run.text[:40]!r}"
                            ))

        # text overflow heuristic
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            line_count = sum(max(1, len(p.text) // 80 + 1) for p in tf.paragraphs)
            try:
                box_h_inches = shape.height / 914400
            except Exception:
                continue
            # rough cap: ~3.5 lines per inch at 18pt
            cap = int(box_h_inches * 3.5) + 2
            if line_count > cap:
                issues.append(Issue("warning", "overflow", idx,
                                    f"text shape may overflow (~{line_count} lines, cap ~{cap})"))

    # duplicate titles
    seen = {}
    for idx, t in titles:
        seen.setdefault(t, []).append(idx)
    for t, idxs in seen.items():
        if len(idxs) > 1:
            issues.append(Issue("warning", "dup-title", idxs[1],
                                f"duplicate title {t!r} on slides {idxs}"))

    # title capitalization consistency
    if titles:
        title_case = sum(1 for _, t in titles if _is_title_case(t))
        sentence_case = sum(1 for _, t in titles if _is_sentence_case(t))
        total = len(titles)
        if 0 < title_case < total and 0 < sentence_case < total:
            issues.append(Issue("warning", "caps-mixed", 0,
                                f"title casing inconsistent ({title_case} title-case, {sentence_case} sentence-case)"))

    return issues


def _detect_title(slide) -> str | None:
    """Best-effort: largest top-of-slide text run."""
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            top = shape.top or 0
        except Exception:
            top = 0
        text = shape.text_frame.text.strip()
        if text and top < 1500000:  # rough: top half-inch
            candidates.append((top, text))
    if not candidates:
        return None
    candidates.sort()
    return candidates[0][1].splitlines()[0].strip()


_WORD_RE = re.compile(r"\w+")
SMALL_WORDS = {"a", "an", "the", "and", "or", "but", "if", "of", "in", "on", "to", "vs", "for", "by", "with"}


def _is_title_case(s: str) -> bool:
    words = _WORD_RE.findall(s)
    if len(words) < 2:
        return False
    cap = sum(1 for w in words if w[0].isupper() or w.lower() in SMALL_WORDS)
    return cap >= len(words) - 1


def _is_sentence_case(s: str) -> bool:
    words = _WORD_RE.findall(s)
    if not words:
        return False
    if not words[0][0].isupper():
        return False
    rest = words[1:]
    if not rest:
        return False
    upper_rest = sum(1 for w in rest if w[0].isupper())
    # most non-leading words lowercase
    return upper_rest <= max(1, len(rest) // 4)

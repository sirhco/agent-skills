"""Structural + textual diff between two .pptx files.

Reduces each deck to (per-slide) {title, body_text, shape_kinds}.
Returns dict {differs: bool, text: [unified-diff lines], shapes: [...]}.
"""
from __future__ import annotations

import difflib
from pathlib import Path

from pptx import Presentation


def _summary(deck_path: Path):
    prs = Presentation(str(deck_path))
    out = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        body_lines = []
        kinds = []
        for shape in slide.shapes:
            kinds.append(shape.shape_type or 0)
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if not title and (shape.top or 0) < 1500000:
                    title = text.splitlines()[0]
                    rest = text.splitlines()[1:]
                    body_lines.extend(rest)
                else:
                    body_lines.extend(text.splitlines())
        out.append({
            "n": idx,
            "title": title,
            "body": "\n".join(body_lines),
            "kinds": sorted(int(k) for k in kinds),
        })
    return out


def diff(a_path: Path, b_path: Path) -> dict:
    a = _summary(a_path)
    b = _summary(b_path)

    text_a = []
    for s in a:
        text_a.append(f"## slide {s['n']}: {s['title']}")
        text_a.extend(s["body"].splitlines())
    text_b = []
    for s in b:
        text_b.append(f"## slide {s['n']}: {s['title']}")
        text_b.extend(s["body"].splitlines())

    udiff = list(difflib.unified_diff(text_a, text_b, fromfile=str(a_path), tofile=str(b_path), lineterm=""))

    shape_diffs = []
    for sa, sb in zip(a, b):
        if sa["kinds"] != sb["kinds"]:
            shape_diffs.append({"slide": sa["n"], "a": sa["kinds"], "b": sb["kinds"]})
    if len(a) != len(b):
        shape_diffs.append({"slide_count": [len(a), len(b)]})

    differs = bool(udiff or shape_diffs)
    return {"differs": differs, "text": udiff, "shapes": shape_diffs}
